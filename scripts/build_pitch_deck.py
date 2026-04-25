"""
Holdfast Protocol — AI Agent Conference NYC Speaker Pitch Deck
Generator script for python-pptx

Run: python scripts/build_pitch_deck.py
Output: docs/holdfast-pitch-deck-aiagentconf-2026.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# --- Brand colors ---
BG_DARK       = RGBColor(0x0D, 0x11, 0x17)   # #0D1117 — slide background
GREEN         = RGBColor(0x14, 0xF1, 0x95)   # #14F195 — Solana green accent
PURPLE        = RGBColor(0x99, 0x45, 0xFF)   # #9945FF — Solana purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GREY          = RGBColor(0x8B, 0x94, 0x9E)   # muted text
RED_WARN      = RGBColor(0xFF, 0x6B, 0x6B)   # warning/problem accent
YELLOW        = RGBColor(0xFF, 0xD6, 0x60)   # roadmap yellow

# --- Slide dimensions (widescreen 16:9) ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    """Return the blank slide layout."""
    return prs.slide_layouts[6]  # index 6 = blank


def add_bg(slide, color=BG_DARK):
    """Fill the slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    """Add a filled rectangle shape."""
    from pptx.util import Pt as UPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  default_size=Pt(16), default_color=WHITE, default_bold=False,
                  line_spacing=1.15):
    """
    Add a text box with multiple lines.
    lines = list of (text, size, color, bold, italic) or just strings.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, color, bold, italic = item, default_size, default_color, default_bold, False
        else:
            text = item[0]
            size  = item[1] if len(item) > 1 else default_size
            color = item[2] if len(item) > 2 else default_color
            bold  = item[3] if len(item) > 3 else default_bold
            italic= item[4] if len(item) > 4 else False

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txBox


def set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, text="Casemate Labs  ·  holdfastprotocol.com  ·  @CasemateLabs"):
    """Add a footer strip at the bottom of the slide."""
    add_rect(slide,
             left=Inches(0), top=Inches(7.15),
             width=Inches(13.333), height=Inches(0.35),
             fill_color=RGBColor(0x16, 0x1B, 0x22))
    add_text(slide, text,
             left=Inches(0.3), top=Inches(7.17),
             width=Inches(12.7), height=Inches(0.3),
             font_size=Pt(10), color=GREY, align=PP_ALIGN.LEFT)


def add_slide_label(slide, label, number=None):
    """Add a small label/pill in the top-right corner."""
    if number:
        label_text = f"{number}  {label}"
    else:
        label_text = label
    add_text(slide, label_text,
             left=Inches(10.5), top=Inches(0.15),
             width=Inches(2.6), height=Inches(0.35),
             font_size=Pt(11), color=GREY, align=PP_ALIGN.RIGHT)


def add_accent_line(slide, color=GREEN, top=Inches(1.1)):
    """Add a thin horizontal accent line below the title."""
    line = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.333), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================

def slide_01_title(prs):
    """Slide 1 — Title"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)

    # Large background wordmark glyph (decorative — top-right faded)
    add_text(slide, "H", left=Inches(9.5), top=Inches(-0.5), width=Inches(4), height=Inches(5),
             font_size=Pt(280), color=RGBColor(0x16, 0x1B, 0x22), bold=True, align=PP_ALIGN.RIGHT)

    # Green accent bar left edge
    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=GREEN)

    # Main title
    add_text(slide, "Holdfast Protocol",
             left=Inches(0.7), top=Inches(1.8), width=Inches(10), height=Inches(1.4),
             font_size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tagline
    add_text(slide, "Trust infrastructure for the AI agent economy.",
             left=Inches(0.7), top=Inches(3.2), width=Inches(10), height=Inches(0.7),
             font_size=Pt(26), bold=False, color=GREEN, align=PP_ALIGN.LEFT)

    # Divider
    add_accent_line(slide, color=GREY, top=Inches(4.1))

    # Conference info
    add_text(slide, "AI Agent Conference NYC  ·  May 4–5, 2026  ·  Agentic Engineering Track",
             left=Inches(0.7), top=Inches(4.25), width=Inches(10), height=Inches(0.5),
             font_size=Pt(16), color=GREY, align=PP_ALIGN.LEFT)

    add_text(slide, "Casemate Labs  ·  holdfastprotocol.com",
             left=Inches(0.7), top=Inches(4.75), width=Inches(8), height=Inches(0.4),
             font_size=Pt(14), color=GREY, align=PP_ALIGN.LEFT)

    set_notes(slide, """Open with a beat of silence. Let the slide breathe.

"Today I'm going to show you the missing accountability layer for the agent economy. Five minutes. Three primitives. One devnet you can use today."

Pace: slow, confident. Stay on this slide for 5-8 seconds before advancing. It sets the register for the whole talk.""")
    return slide


def slide_02_problem(prs):
    """Slide 2 — The Problem"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "The Problem", 2)

    # Title
    add_text(slide, "Agents are moving real money.",
             left=Inches(0.6), top=Inches(0.5), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)
    add_text(slide, "There is no standard for what happens when they go wrong.",
             left=Inches(0.6), top=Inches(1.2), width=Inches(12), height=Inches(0.6),
             font_size=Pt(28), bold=False, color=RED_WARN)

    add_accent_line(slide, color=RED_WARN, top=Inches(1.95))

    # Bullet points
    bullets = [
        ("AI agents executing transactions autonomously", Pt(20), WHITE, False),
        ("Signing transactions with software keys in .env files", Pt(20), WHITE, False),
        ("No verifiable identity — a key proves nothing", Pt(20), WHITE, False),
        ("No on-chain track record of past behavior", Pt(20), WHITE, False),
        ("No enforceable consequences for failure or fraud", Pt(20), RED_WARN, True),
    ]
    y = Inches(2.15)
    for text, size, color, bold in bullets:
        add_text(slide, f"  ▸  {text}",
                 left=Inches(0.8), top=y, width=Inches(11), height=Inches(0.45),
                 font_size=size, color=color, bold=bold)
        y += Inches(0.55)

    # Right side — stat callout
    add_rect(slide, left=Inches(8.8), top=Inches(2.2), width=Inches(4.1), height=Inches(2.8),
             fill_color=RGBColor(0x16, 0x1B, 0x22),
             line_color=RED_WARN, line_width=Pt(1))
    add_text(slide, "21,000+",
             left=Inches(8.9), top=Inches(2.4), width=Inches(3.9), height=Inches(1.0),
             font_size=Pt(52), bold=True, color=RED_WARN, align=PP_ALIGN.CENTER)
    add_text(slide, "AI agents deployed in the Colosseum\nhackathon alone.\n\nNo accountability layer.",
             left=Inches(8.9), top=Inches(3.3), width=Inches(3.9), height=Inches(1.5),
             font_size=Pt(14), color=GREY, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: 0:00–0:45

"AI agents are executing transactions, managing funds, and settling contracts without human intervention right now. The Colosseum AI hackathon alone produced 21,000 agents in February. These agents are making real decisions with real money."

"But the infrastructure they run on was never designed for this. Agents authenticate with software keys stored in environment variables. One compromised service and your agent's signing key is gone."

"There's no standard for 'who is this agent', 'can it be trusted', or 'what happens if it cheats.'"

"That's not a small gap. That's the missing accountability layer for the entire agent economy." """)
    return slide


def slide_03_trust_gap(prs):
    """Slide 3 — The Trust Gap"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "The Trust Gap", 3)

    add_text(slide, "The trust gap has three parts:",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.65),
             font_size=Pt(36), bold=True, color=WHITE)
    add_accent_line(slide, color=GREEN, top=Inches(1.1))

    # Three cards
    cards = [
        ("IDENTITY",   "Who is this\nagent, really?",   "Software keys\ncan be copied.",   GREEN),
        ("REPUTATION", "Has this agent\nbehaved before?", "No on-chain\ntrack record.",     PURPLE),
        ("ENFORCEMENT","What happens\nwhen it fails?",   "No programmable\nsettlement.",    RED_WARN),
    ]
    card_w = Inches(3.8)
    card_h = Inches(4.5)
    card_top = Inches(1.3)
    for i, (title, q, gap, color) in enumerate(cards):
        x = Inches(0.5 + i * 4.3)
        add_rect(slide, left=x, top=card_top, width=card_w, height=card_h,
                 fill_color=RGBColor(0x16, 0x1B, 0x22),
                 line_color=color, line_width=Pt(1.5))
        # Card accent top bar
        add_rect(slide, left=x, top=card_top, width=card_w, height=Inches(0.07),
                 fill_color=color)
        add_text(slide, title,
                 left=x + Inches(0.15), top=card_top + Inches(0.2),
                 width=card_w - Inches(0.3), height=Inches(0.5),
                 font_size=Pt(18), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, q,
                 left=x + Inches(0.15), top=card_top + Inches(0.85),
                 width=card_w - Inches(0.3), height=Inches(1.2),
                 font_size=Pt(22), bold=False, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, gap,
                 left=x + Inches(0.15), top=card_top + Inches(2.8),
                 width=card_w - Inches(0.3), height=Inches(1.2),
                 font_size=Pt(17), bold=False, color=GREY, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: 0:45–1:20

"The trust gap is actually three problems layered on top of each other."

"First: identity. A software key tells you nothing — it can be copied, stolen, or spoofed."

"Second: reputation. There's no on-chain track record. You can't know whether the agent you're about to trust has ever successfully completed anything."

"Third: enforcement. Even if you could verify identity and reputation, there's no programmable settlement layer — no way to lock funds, define release conditions, and resolve disputes with on-chain finality."

"Holdfast Protocol solves all three." """)
    return slide


def slide_04_intro(prs):
    """Slide 4 — Introducing Holdfast Protocol"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Holdfast Protocol", 4)

    add_text(slide, "Holdfast Protocol",
             left=Inches(0.6), top=Inches(0.4), width=Inches(10), height=Inches(0.75),
             font_size=Pt(44), bold=True, color=WHITE)
    add_text(slide, "Three composable primitives for AI agent trust. Native on Solana.",
             left=Inches(0.6), top=Inches(1.1), width=Inches(12), height=Inches(0.5),
             font_size=Pt(22), color=GREEN)
    add_accent_line(slide, color=GREEN, top=Inches(1.7))

    # Three primitive boxes with arrows between
    primitives = [
        ("Identity",    "Hardware-attested\nagent identities",  GREEN),
        ("Reputation",  "On-chain oracle\n(CPI-readable)",     PURPLE),
        ("Escrow",      "Programmable\nsettlement",             RGBColor(0xFF, 0xA5, 0x00)),
    ]
    box_w = Inches(3.4)
    box_h = Inches(2.8)
    box_top = Inches(2.0)
    for i, (name, desc, color) in enumerate(primitives):
        x = Inches(0.6 + i * 4.2)
        add_rect(slide, left=x, top=box_top, width=box_w, height=box_h,
                 fill_color=RGBColor(0x16, 0x1B, 0x22),
                 line_color=color, line_width=Pt(2))
        add_rect(slide, left=x, top=box_top, width=box_w, height=Inches(0.07),
                 fill_color=color)
        add_text(slide, name,
                 left=x + Inches(0.1), top=box_top + Inches(0.3),
                 width=box_w - Inches(0.2), height=Inches(0.65),
                 font_size=Pt(30), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 left=x + Inches(0.1), top=box_top + Inches(1.1),
                 width=box_w - Inches(0.2), height=Inches(1.4),
                 font_size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)
        # Arrow between boxes
        if i < 2:
            add_text(slide, "→",
                     left=x + box_w + Inches(0.1), top=box_top + Inches(1.0),
                     width=Inches(0.6), height=Inches(0.6),
                     font_size=Pt(28), color=GREY, align=PP_ALIGN.CENTER)

    # Bottom badge
    add_rect(slide, left=Inches(4.5), top=Inches(5.1), width=Inches(4.3), height=Inches(0.55),
             fill_color=RGBColor(0x10, 0x2A, 0x1E), line_color=GREEN, line_width=Pt(1))
    add_text(slide, "Devnet live  ·  npm install @holdfastprotocol/sdk@devnet",
             left=Inches(4.5), top=Inches(5.15), width=Inches(4.3), height=Inches(0.45),
             font_size=Pt(13), color=GREEN, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: 1:20–1:45

"Holdfast Protocol is trust infrastructure for the AI agent economy, deployed natively on Solana. It gives autonomous agents three things they've never had: a verifiable identity, an on-chain reputation score, and a programmable escrow contract they cannot break."

"Each primitive is independent but composable. You can read an agent's reputation before accepting any pact. You can gate escrow creation on a minimum reputation score. The chain enforces it."

"All three are live on devnet today."

Advance quickly — you have 5 minutes total.""")
    return slide


def slide_05_identity(prs):
    """Slide 5 — Primitive 1: Identity"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Identity", 5)

    # Header
    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=GREEN)
    add_text(slide, "Primitive 1",
             left=Inches(0.5), top=Inches(0.3), width=Inches(6), height=Inches(0.45),
             font_size=Pt(16), color=GREEN, bold=True)
    add_text(slide, "Identity",
             left=Inches(0.5), top=Inches(0.7), width=Inches(8), height=Inches(0.75),
             font_size=Pt(44), bold=True, color=WHITE)
    add_text(slide, "Hardware-attested agent identities on Solana.",
             left=Inches(0.5), top=Inches(1.4), width=Inches(10), height=Inches(0.5),
             font_size=Pt(20), color=GREEN)
    add_accent_line(slide, color=GREEN, top=Inches(2.0))

    # Flow steps
    steps = [
        ("Agent generates secp256r1 (P-256) keypair", WHITE),
        ("Same primitive: hardware security keys, Apple Secure Enclave, WebAuthn", GREY),
        ("Holdfast Protocol anchors this key on-chain via Solana's native secp256r1 precompile (SIMD-48)", WHITE),
        ("AgentWallet PDA — permanently on-chain, verifiable by anyone", GREEN),
    ]
    y = Inches(2.2)
    for i, (text, color) in enumerate(steps):
        add_rect(slide, left=Inches(0.5), top=y + Inches(0.12),
                 width=Inches(0.35), height=Inches(0.35),
                 fill_color=GREEN if color == GREEN else RGBColor(0x16, 0x1B, 0x22),
                 line_color=GREEN, line_width=Pt(1))
        add_text(slide, str(i + 1),
                 left=Inches(0.52), top=y + Inches(0.1), width=Inches(0.32), height=Inches(0.38),
                 font_size=Pt(13), bold=True, color=WHITE if color == GREEN else GREEN,
                 align=PP_ALIGN.CENTER)
        add_text(slide, text,
                 left=Inches(1.05), top=y, width=Inches(11.5), height=Inches(0.6),
                 font_size=Pt(18), color=color)
        if i < len(steps) - 1:
            add_text(slide, "↓", left=Inches(0.58), top=y + Inches(0.5),
                     width=Inches(0.25), height=Inches(0.35), font_size=Pt(14), color=GREY,
                     align=PP_ALIGN.CENTER)
        y += Inches(0.85)

    # Quote
    add_rect(slide, left=Inches(0.5), top=Inches(5.8), width=Inches(12.3), height=Inches(0.7),
             fill_color=RGBColor(0x10, 0x2A, 0x1E), line_color=GREEN, line_width=Pt(1))
    add_text(slide, '"Not a software key in an .env file.  A hardware-bound identity on Solana."',
             left=Inches(0.7), top=Inches(5.88), width=Inches(12.0), height=Inches(0.55),
             font_size=Pt(16), color=GREEN, italic=True, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: ~1:45

"The identity primitive uses secp256r1 — P-256 — the same elliptic curve used by hardware security keys, Apple Secure Enclave, and WebAuthn. Hardware-rootable. Can't be trivially copied."

"When an agent registers, it generates a P-256 keypair and proves possession on-chain using Solana's native secp256r1 precompile — SIMD-48. Instruction-level primitive, not a smart contract. Attestation verified in the same instruction as program execution. No separate round-trip."

"Full TPM and TEE hardware attestation — integrating with our Hardline Protocol — is on the roadmap. What's live today is secp256r1 self-attestation. Already meaningfully stronger than any software key."

TIME CHECK: you should be at ~1:20 on the clock.""")
    return slide


def slide_06_reputation(prs):
    """Slide 6 — Primitive 2: Reputation"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Reputation", 6)

    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=PURPLE)
    add_text(slide, "Primitive 2",
             left=Inches(0.5), top=Inches(0.3), width=Inches(6), height=Inches(0.45),
             font_size=Pt(16), color=PURPLE, bold=True)
    add_text(slide, "Reputation",
             left=Inches(0.5), top=Inches(0.7), width=Inches(8), height=Inches(0.75),
             font_size=Pt(44), bold=True, color=WHITE)
    add_text(slide, "On-chain reputation oracle.",
             left=Inches(0.5), top=Inches(1.4), width=Inches(10), height=Inches(0.5),
             font_size=Pt(20), color=PURPLE)
    add_accent_line(slide, color=PURPLE, top=Inches(2.0))

    # Left column — description
    left_items = [
        ("Score: 0 – 10,000 basis points  (5,000 = neutral)", WHITE),
        ("Updated by: the protocol escrow program (not a human)", GREY),
        ("Every pact outcome posted on-chain: fulfilled or disputed", WHITE),
        ("Score decays lazily toward neutral when inactive", GREY),
        ("Readable by any Solana program via CPI — one account read", PURPLE),
        ("No oracle fee. No bridge. No trust assumption.", WHITE),
    ]
    y = Inches(2.2)
    for text, color in left_items:
        add_text(slide, f"  ▸  {text}",
                 left=Inches(0.5), top=y, width=Inches(7.5), height=Inches(0.42),
                 font_size=Pt(16), color=color)
        y += Inches(0.48)

    # Right column — code snippet
    code_bg = RGBColor(0x16, 0x1B, 0x22)
    add_rect(slide, left=Inches(8.3), top=Inches(2.1), width=Inches(4.7), height=Inches(3.0),
             fill_color=code_bg, line_color=PURPLE, line_width=Pt(1))
    code_lines = [
        ("const ok = await client", Pt(12), GREY, False),
        ("  .reputation", Pt(12), GREY, False),
        ("  .meetsRequirements(", Pt(12), GREY, False),
        ("    agentPubkey, {", Pt(12), WHITE, False),
        ("    minScore: 6000,", Pt(12), GREEN, False),
        ("    minTier: VerifTier.Attested,", Pt(12), GREEN, False),
        ("    minPacts: 3,", Pt(12), GREEN, False),
        ("  }", Pt(12), WHITE, False),
        (");", Pt(12), WHITE, False),
    ]
    cy = Inches(2.25)
    for text, size, color, bold in code_lines:
        add_text(slide, text,
                 left=Inches(8.45), top=cy, width=Inches(4.4), height=Inches(0.3),
                 font_size=size, color=color, bold=bold)
        cy += Inches(0.3)

    add_text(slide, "Pre-flight matches on-chain constraint.\nReturns false (not throws) for unknown agents.",
             left=Inches(8.3), top=Inches(5.2), width=Inches(4.7), height=Inches(0.7),
             font_size=Pt(12), color=GREY, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: ~2:00

"The reputation primitive is an on-chain oracle. Scores run 0 to 10,000 basis points — 5,000 is neutral. Every time an agent fulfills or disputes a pact, the escrow settlement program — a program PDA, not a human — posts a signed reputation update."

"No centralized reputation manipulation. The oracle authority is a program account. The whole history is on-chain, auditable by anyone."

"The killer feature: because the reputation account is a PDA, any other Solana program can read it via CPI. One account read. No oracle fee, no bridge, no cross-chain message."

"Scores decay lazily toward neutral when inactive — a simple, transparent design to prevent permanent reputation from stale accounts." """)
    return slide


def slide_07_escrow(prs):
    """Slide 7 — Primitive 3: Escrow"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Escrow", 7)

    ORANGE = RGBColor(0xFF, 0xA5, 0x00)
    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=ORANGE)
    add_text(slide, "Primitive 3",
             left=Inches(0.5), top=Inches(0.3), width=Inches(6), height=Inches(0.45),
             font_size=Pt(16), color=ORANGE, bold=True)
    add_text(slide, "Escrow",
             left=Inches(0.5), top=Inches(0.7), width=Inches(8), height=Inches(0.75),
             font_size=Pt(44), bold=True, color=WHITE)
    add_text(slide, "Programmable settlement for AI agent commerce.",
             left=Inches(0.5), top=Inches(1.4), width=Inches(10), height=Inches(0.5),
             font_size=Pt(20), color=ORANGE)
    add_accent_line(slide, color=ORANGE, top=Inches(2.0))

    # Lifecycle flow
    stages = ["Create Pact", "Fund", "Release", "Dispute\n(optional)", "Resolve"]
    stage_colors = [GREEN, GREEN, GREEN, RED_WARN, ORANGE]
    sx = Inches(0.5)
    sy = Inches(2.3)
    box_w = Inches(2.1)
    box_h = Inches(0.8)
    for i, (s, c) in enumerate(zip(stages, stage_colors)):
        add_rect(slide, left=sx, top=sy, width=box_w, height=box_h,
                 fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=c, line_width=Pt(1.5))
        add_text(slide, s, left=sx, top=sy, width=box_w, height=box_h,
                 font_size=Pt(14), bold=True, color=c, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_text(slide, "→",
                     left=sx + box_w + Inches(0.05), top=sy + Inches(0.2),
                     width=Inches(0.35), height=Inches(0.4),
                     font_size=Pt(18), color=GREY, align=PP_ALIGN.CENTER)
        sx += Inches(2.5)

    # Details
    details = [
        ("Task-based, milestone-gated, or time-locked release conditions.", WHITE),
        ("Funds lock at pact initiation.", WHITE),
        ("Release opens a 7-day dispute window before final settlement.", WHITE),
        ("Disputes trigger arbiter resolution with on-chain finality.", WHITE),
    ]
    y = Inches(3.4)
    for text, color in details:
        add_text(slide, f"  ▸  {text}",
                 left=Inches(0.5), top=y, width=Inches(8.5), height=Inches(0.42),
                 font_size=Pt(17), color=color)
        y += Inches(0.5)

    # Reputation gating callout
    add_rect(slide, left=Inches(9.0), top=Inches(2.25), width=Inches(4.0), height=Inches(2.8),
             fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=ORANGE, line_width=Pt(1.5))
    add_text(slide, "Reputation gating built in:",
             left=Inches(9.1), top=Inches(2.4), width=Inches(3.8), height=Inches(0.45),
             font_size=Pt(14), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    add_text(slide, "Pre-flight checks counterparty\nscore before any funds move.\n\nIf threshold not met →\nReputationThresholdNotMet\n(before fees, before commitment)",
             left=Inches(9.1), top=Inches(2.9), width=Inches(3.8), height=Inches(1.8),
             font_size=Pt(13), color=WHITE, align=PP_ALIGN.LEFT)

    # Quote
    add_rect(slide, left=Inches(0.5), top=Inches(5.9), width=Inches(12.3), height=Inches(0.65),
             fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=ORANGE, line_width=Pt(1))
    add_text(slide, '"The agent can\'t just walk away.  The chain holds it accountable."',
             left=Inches(0.7), top=Inches(5.98), width=Inches(12.0), height=Inches(0.5),
             font_size=Pt(16), color=ORANGE, italic=True, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: ~2:30

"The escrow primitive is programmable settlement. You define release conditions at creation time: task-based, milestone-gated, or time-locked. Funds lock at pact initiation. Releasing opens a dispute window before final settlement."

"The three primitives compose naturally. When you create a pact, you can require a minimum reputation threshold for the counterparty. The SDK runs a pre-flight before the transaction hits the chain — if the counterparty doesn't qualify, clean error before any fees are paid."

"The dispute path has on-chain finality. An arbiter — a multisig or another program — resolves disputes and the outcome posts back to the reputation oracle. Bad behavior has real, on-chain consequences."

TIME CHECK: you should be at ~2:30 on the clock.""")
    return slide


def slide_08_why_solana(prs):
    """Slide 8 — Why Solana"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Why Solana", 8)

    add_text(slide, "Why Solana — not EVM?",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)
    add_accent_line(slide, color=GREEN, top=Inches(1.2))

    # Comparison table
    headers = ["", "Holdfast Protocol\non Solana", "EAS / EVM forks"]
    rows = [
        ("Attestation",        "Same instruction as\nprogram execution",   "Separate contract call"),
        ("Cost per operation", "Sub-cent",                                  "EVM gas — compounds at scale"),
        ("Finality",           "~400ms",                                    "~12 seconds (Ethereum)"),
        ("Composability",      "Native CPI — any program reads reputation","Bridge required for non-EVM agents"),
        ("AI agent focus",     "Built for agent commerce",                  "General purpose"),
    ]
    col_w = [Inches(3.0), Inches(4.5), Inches(4.5)]
    row_h = Inches(0.7)
    table_top = Inches(1.35)
    table_left = Inches(0.6)

    # Header row
    x = table_left
    for i, h in enumerate(headers):
        bg = RGBColor(0x10, 0x2A, 0x1E) if i == 1 else RGBColor(0x16, 0x1B, 0x22)
        add_rect(slide, left=x, top=table_top, width=col_w[i], height=row_h,
                 fill_color=bg)
        add_text(slide, h, left=x + Inches(0.1), top=table_top + Inches(0.1),
                 width=col_w[i] - Inches(0.2), height=row_h - Inches(0.15),
                 font_size=Pt(14), bold=True,
                 color=GREEN if i == 1 else WHITE,
                 align=PP_ALIGN.CENTER)
        x += col_w[i]

    # Data rows
    for r, (label, sol_val, evm_val) in enumerate(rows):
        y = table_top + row_h * (r + 1)
        vals = [label, sol_val, evm_val]
        x = table_left
        for i, v in enumerate(vals):
            bg = RGBColor(0x0F, 0x24, 0x18) if i == 1 else RGBColor(0x13, 0x17, 0x1E)
            add_rect(slide, left=x, top=y, width=col_w[i], height=row_h,
                     fill_color=bg,
                     line_color=RGBColor(0x21, 0x26, 0x2D), line_width=Pt(0.5))
            add_text(slide, v, left=x + Inches(0.1), top=y + Inches(0.05),
                     width=col_w[i] - Inches(0.2), height=row_h - Inches(0.1),
                     font_size=Pt(13),
                     color=GREEN if i == 1 else (WHITE if i == 0 else GREY),
                     bold=(i == 0),
                     align=PP_ALIGN.LEFT)
            x += col_w[i]

    # Key stat callouts at bottom
    stats = [
        ("~400ms", "finality"),
        ("<$0.001", "per operation"),
        ("Native secp256r1", "precompile (SIMD-48)"),
    ]
    bx = Inches(0.6)
    for val, label in stats:
        add_rect(slide, left=bx, top=Inches(5.85), width=Inches(3.8), height=Inches(0.85),
                 fill_color=RGBColor(0x10, 0x2A, 0x1E), line_color=GREEN, line_width=Pt(1))
        add_text(slide, val,
                 left=bx + Inches(0.1), top=Inches(5.88), width=Inches(3.6), height=Inches(0.4),
                 font_size=Pt(22), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, label,
                 left=bx + Inches(0.1), top=Inches(6.25), width=Inches(3.6), height=Inches(0.35),
                 font_size=Pt(13), color=GREY, align=PP_ALIGN.CENTER)
        bx += Inches(4.1)

    set_notes(slide, """TIME: ~3:00

"Why Solana? Two reasons: cost and composability."

"EAS and its EVM derivatives work — but the design shows it. Attestation records are separate from the execution environment. Verification requires cross-contract calls with EVM gas overhead. At 12-second block times and real gas costs, that's a real constraint for agent economies."

"Holdfast Protocol is Solana-native. The secp256r1 precompile is an instruction-level primitive — attestation in the same instruction as the program call. ~400ms finality at a fraction of a cent. For agents signing dozens of transactions per hour, that cost difference compounds."

"And composability: any Solana program can read agent reputation via a single CPI account read. No oracle fee, no bridge, no cross-chain message." """)
    return slide


def slide_09_demo(prs):
    """Slide 9 — Live Demo"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Live Demo", 9)

    add_text(slide, "Live Devnet Demo",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)
    add_accent_line(slide, color=GREEN, top=Inches(1.2))

    # Stages
    stages = [
        ("Stage 1", "Register agent wallet",
         "Generate secp256r1 keypair →\nregisterAgentWallet tx confirmed on Solana devnet", GREEN),
        ("Stage 2", "Oracle reputation update",
         "initReputation (score: 5000 bp, neutral) →\nupdateReputation +200bp → Score: 5200/10000 on-chain", PURPLE),
        ("Stage 3", "SDK attestation query",
         "client.reputation.get(agentPubkey) →\nClean attestation report: score, tier, history", RGBColor(0xFF, 0xA5, 0x00)),
    ]
    y = Inches(1.4)
    for label, title, desc, color in stages:
        add_rect(slide, left=Inches(0.6), top=y, width=Inches(8.5), height=Inches(1.4),
                 fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=color, line_width=Pt(1.5))
        add_rect(slide, left=Inches(0.6), top=y, width=Inches(0.07), height=Inches(1.4),
                 fill_color=color)
        add_text(slide, label,
                 left=Inches(0.85), top=y + Inches(0.1), width=Inches(2), height=Inches(0.35),
                 font_size=Pt(12), bold=True, color=color)
        add_text(slide, title,
                 left=Inches(0.85), top=y + Inches(0.4), width=Inches(8), height=Inches(0.4),
                 font_size=Pt(18), bold=True, color=WHITE)
        add_text(slide, desc,
                 left=Inches(0.85), top=y + Inches(0.78), width=Inches(8), height=Inches(0.5),
                 font_size=Pt(13), color=GREY)
        y += Inches(1.55)

    # Demo video placeholder
    add_rect(slide, left=Inches(9.4), top=Inches(1.35), width=Inches(3.5), height=Inches(4.0),
             fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=GREEN, line_width=Pt(1.5))
    add_text(slide, "▶",
             left=Inches(9.4), top=Inches(2.2), width=Inches(3.5), height=Inches(1.0),
             font_size=Pt(52), color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "WATCH DEMO\n4 min",
             left=Inches(9.4), top=Inches(3.1), width=Inches(3.5), height=Inches(0.7),
             font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "[ link to be added\nby Video Editor — Apr 28 ]",
             left=Inches(9.4), top=Inches(3.75), width=Inches(3.5), height=Inches(0.6),
             font_size=Pt(11), color=GREY, align=PP_ALIGN.CENTER)

    # Disclaimer
    add_rect(slide, left=Inches(0.6), top=Inches(6.25), width=Inches(12.3), height=Inches(0.45),
             fill_color=RGBColor(0x2A, 0x16, 0x10), line_color=RED_WARN, line_width=Pt(1))
    add_text(slide, "⚠  Pre-audit software. Deployed on Solana devnet. Not for mainnet or production use.",
             left=Inches(0.8), top=Inches(6.31), width=Inches(12.0), height=Inches(0.35),
             font_size=Pt(12), color=RED_WARN, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: 3:15–4:15

"Let me show you what this looks like in practice. This is a real Solana devnet run — not a simulation."

[If showing live demo video]: "We're watching a 90-second clip. Three moments: the registerAgentWallet transaction landing in a real Solana block, the oracle reputation update with score ticking from 5000 to 5200, and the SDK query returning a clean attestation report."

[If not playing video]: "The full demo is at the QR code — I'd encourage you to run it yourself. ~5 minutes from npm install to your first on-chain agent registration."

"Everything you see is live on devnet. Program ID is in Solana Explorer. The SDK is on npm right now."

NOTE: Replace [link placeholder] with Video Editor's recording link once received (expected Apr 28). """)
    return slide


def slide_10_current_state(prs):
    """Slide 10 — Current State"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Current State", 10)

    add_text(slide, "What's live today",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)
    add_text(slide, "No overclaiming. Here's exactly where we are.",
             left=Inches(0.6), top=Inches(1.1), width=Inches(10), height=Inches(0.45),
             font_size=Pt(18), color=GREY)
    add_accent_line(slide, color=GREEN, top=Inches(1.65))

    # Status items
    live_items = [
        "Reputation read/write on devnet",
        "Escrow create / fund / release on devnet",
        "secp256r1 self-attestation (localnet; devnet pending cluster upgrade)",
        "@holdfastprotocol/sdk@0.1.0-devnet.1 published to npm",
        "Operator dashboard (devnet) — reputation, escrow, custody views",
        "Off-chain indexer deployed (indexer.devnet.holdfastprotocol.com)",
    ]
    roadmap_items = [
        "Hardware TPM / TEE attestation (Hardline cross-CPI) — Q4 2026",
        "Mainnet deployment — after external security audit (in progress)",
        "Protocol fees on production usage",
    ]

    y = Inches(1.8)
    for item in live_items:
        add_text(slide, f"  ✅  {item}",
                 left=Inches(0.6), top=y, width=Inches(9.5), height=Inches(0.4),
                 font_size=Pt(16), color=WHITE)
        y += Inches(0.46)

    y += Inches(0.1)
    for item in roadmap_items:
        add_text(slide, f"  🔜  {item}",
                 left=Inches(0.6), top=y, width=Inches(9.5), height=Inches(0.4),
                 font_size=Pt(16), color=GREY)
        y += Inches(0.46)

    # No-token callout
    add_rect(slide, left=Inches(9.8), top=Inches(1.8), width=Inches(3.1), height=Inches(2.5),
             fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=GREEN, line_width=Pt(1.5))
    add_text(slide, "No token.",
             left=Inches(9.9), top=Inches(2.0), width=Inches(2.9), height=Inches(0.55),
             font_size=Pt(26), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Revenue: protocol fees\non real usage.\n\nRegistrations,\nescrow settlements.",
             left=Inches(9.9), top=Inches(2.55), width=Inches(2.9), height=Inches(1.5),
             font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: ~4:15 (brief slide — 20 seconds)

"Here's exactly what's live and what's not — no overclaiming."

"The SDK is published. The programs are on devnet. The operator dashboard is running. You can install, connect a devnet wallet, register an agent identity, and read your reputation score today."

"What's not live: full hardware TPM/TEE attestation — roadmap for Q4. Mainnet deployment is gated on the external security audit, which is in progress."

"No token. Business model is protocol fees on real usage."

"We're telling you exactly what's done and what isn't. That's the kind of team you want to build on." """)
    return slide


def slide_11_competitive(prs):
    """Slide 11 — Competitive Differentiation"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Differentiation", 11)

    add_text(slide, "How we're different",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)
    add_accent_line(slide, color=GREEN, top=Inches(1.2))

    # Comparison table
    cols = ["", "Holdfast", "EAS / EVM", "Vouch", "Warden"]
    col_w = [Inches(2.8), Inches(2.5), Inches(2.3), Inches(2.3), Inches(2.3)]
    rows = [
        ("On-chain enforcement",    "✅ Solana",  "✅ EVM",    "❌",    "Partial"),
        ("Hardware attestation",    "✅ P-256\nTPM roadmap", "❌",  "❌ SSH/Ed25519", "✅ MPC"),
        ("Reputation oracle",       "✅ CPI-readable", "❌",  "Social",  "❌"),
        ("Programmable escrow",     "✅ + dispute",  "❌",    "❌",    "Partial"),
        ("Solana-native",           "✅",          "❌",      "❌",    "❌"),
        ("AI agent focus",          "✅",          "General", "IdP / DevOps", "Enterprise"),
    ]
    row_h = Inches(0.63)
    table_top = Inches(1.3)
    table_left = Inches(0.3)

    # Header
    x = table_left
    for i, h in enumerate(cols):
        bg = RGBColor(0x10, 0x2A, 0x1E) if i == 1 else RGBColor(0x16, 0x1B, 0x22)
        add_rect(slide, left=x, top=table_top, width=col_w[i], height=row_h,
                 fill_color=bg)
        add_text(slide, h, left=x + Inches(0.05), top=table_top + Inches(0.15),
                 width=col_w[i] - Inches(0.1), height=row_h - Inches(0.2),
                 font_size=Pt(14), bold=True,
                 color=GREEN if i == 1 else WHITE,
                 align=PP_ALIGN.CENTER)
        x += col_w[i]

    # Data
    for r, row in enumerate(rows):
        y = table_top + row_h * (r + 1)
        x = table_left
        for i, v in enumerate(row):
            bg = RGBColor(0x0F, 0x24, 0x18) if i == 1 else RGBColor(0x13, 0x17, 0x1E)
            add_rect(slide, left=x, top=y, width=col_w[i], height=row_h,
                     fill_color=bg,
                     line_color=RGBColor(0x21, 0x26, 0x2D), line_width=Pt(0.5))
            is_check = "✅" in v
            is_cross  = v.strip() == "❌"
            color = GREEN if (i == 1 or is_check) else (RED_WARN if is_cross else GREY)
            add_text(slide, v, left=x + Inches(0.05), top=y + Inches(0.1),
                     width=col_w[i] - Inches(0.1), height=row_h - Inches(0.15),
                     font_size=Pt(13),
                     color=color, bold=(i == 0),
                     align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x += col_w[i]

    # Quote
    add_rect(slide, left=Inches(0.3), top=Inches(5.85), width=Inches(12.7), height=Inches(0.65),
             fill_color=RGBColor(0x10, 0x2A, 0x1E), line_color=GREEN, line_width=Pt(1))
    add_text(slide, '"They built the badge.  We built the vault."',
             left=Inches(0.5), top=Inches(5.93), width=Inches(12.3), height=Inches(0.5),
             font_size=Pt(22), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: ~4:15 (30 seconds)

"Three names you'll hear compared to us: EAS, Vouch, Warden."

"EAS: general-purpose attestation on EVM. Not designed for AI agents, not on Solana, no escrow or financial enforcement."

"Vouch: off-chain identity tool — Python library, SSH keys, Git signing, DIDs. Gives an agent a verified ID badge. Cannot hold assets, slash stake, or resolve a dispute. We see Vouch as a potential integration partner — a Vouch DID could feed a Holdfast agent registration. If you're talking to their team here, tell them we want to connect."

"Warden: enterprise MPC key management. Not designed for on-chain agent trust."

"Our moat: Solana-native hardware attestation, on-chain reputation, and programmable escrow in a single composable stack." """)
    return slide


def slide_12_ecosystem(prs):
    """Slide 12 — Ecosystem Fit"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Ecosystem Fit", 12)

    add_text(slide, "Built for the Solana AI agent ecosystem",
             left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.75),
             font_size=Pt(36), bold=True, color=WHITE)
    add_accent_line(slide, color=GREEN, top=Inches(1.2))

    # Three columns
    categories = [
        ("Agent Frameworks", GREEN, [
            ("ElizaOS / ai16z",       "Holdfast plugin: agent\nidentity + escrow"),
            ("Solana Agent Kit",      "Native Holdfast\nProtocol actions"),
            ("Olas / Autonolas",      "Attestation for\nautonomous services"),
        ]),
        ("DeFi Protocols", PURPLE, [
            ("Jupiter",               "Agent-gated routing limits\n(reputation check)"),
            ("Drift Protocol",        "Agent identity + risk\nattestation for perps"),
            ("Jito",                  "Validator-side agent\ntrust layer"),
        ]),
        ("Integrate now", RGBColor(0xFF, 0xA5, 0x00), [
            ("npm install", "@holdfastprotocol\n/sdk@devnet"),
            ("Docs",        "holdfastprotocol.com\n/docs"),
            ("Contact",     "Find us at the\nconference today"),
        ]),
    ]

    col_x = [Inches(0.5), Inches(4.7), Inches(8.9)]
    col_w = Inches(3.9)
    for col_i, (cat_name, col_color, items) in enumerate(categories):
        x = col_x[col_i]
        # Category header
        add_rect(slide, left=x, top=Inches(1.35), width=col_w, height=Inches(0.5),
                 fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=col_color, line_width=Pt(1.5))
        add_text(slide, cat_name, left=x + Inches(0.1), top=Inches(1.38),
                 width=col_w - Inches(0.2), height=Inches(0.42),
                 font_size=Pt(16), bold=True, color=col_color, align=PP_ALIGN.CENTER)

        # Items
        iy = Inches(1.95)
        for item_name, item_desc in items:
            add_rect(slide, left=x, top=iy, width=col_w, height=Inches(1.4),
                     fill_color=RGBColor(0x16, 0x1B, 0x22),
                     line_color=RGBColor(0x21, 0x26, 0x2D), line_width=Pt(0.5))
            add_text(slide, item_name, left=x + Inches(0.15), top=iy + Inches(0.1),
                     width=col_w - Inches(0.3), height=Inches(0.45),
                     font_size=Pt(15), bold=True, color=WHITE)
            add_text(slide, item_desc, left=x + Inches(0.15), top=iy + Inches(0.55),
                     width=col_w - Inches(0.3), height=Inches(0.7),
                     font_size=Pt(13), color=GREY)
            iy += Inches(1.5)

    set_notes(slide, """TIME: ~4:30 (30 seconds)

"We're building for the Solana AI agent ecosystem, not against it."

"Integration proposals in progress with ElizaOS and Solana Agent Kit — any agent on those frameworks can register with Holdfast Protocol in a single SDK call. We want to talk to those teams here this week."

"For DeFi protocols: if you're allowing agents to execute trades or manage positions, the reputation oracle is a CPI call away. Gate agent access to high-risk operations on a minimum reputation score before any funds move."

"We're here specifically to meet builders. If you're in agent frameworks or DeFi and you need an accountability layer, find us." """)
    return slide


def slide_13_cta(prs):
    """Slide 13 — Call to Action"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_footer(slide)
    add_slide_label(slide, "Call to Action", 13)

    # Green left bar
    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=GREEN)

    add_text(slide, "Build with Holdfast Protocol today.",
             left=Inches(0.5), top=Inches(0.5), width=Inches(12), height=Inches(0.75),
             font_size=Pt(40), bold=True, color=WHITE)

    # npm install block
    add_rect(slide, left=Inches(0.5), top=Inches(1.35), width=Inches(12.3), height=Inches(0.65),
             fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=GREEN, line_width=Pt(1.5))
    add_text(slide, "npm install @holdfastprotocol/sdk@devnet",
             left=Inches(0.7), top=Inches(1.43), width=Inches(12.0), height=Inches(0.5),
             font_size=Pt(22), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_accent_line(slide, color=GREY, top=Inches(2.15))

    # Three CTA cards
    ctas = [
        ("1.  Try the SDK", GREEN,
         "devnet · no token · no mainnet risk\n\nnpm install @holdfastprotocol/sdk@devnet\nholdfastprotocol.com/docs"),
        ("2.  Partner with us", PURPLE,
         "ElizaOS · Solana Agent Kit · DeFi protocols\n\nTalk to us today\nat the conference"),
        ("3.  Follow for mainnet", RGBColor(0xFF, 0xA5, 0x00),
         "Audit in progress · Mainnet post-audit\n\n@CasemateLabs\nfor timeline updates"),
    ]
    cx = Inches(0.5)
    for title, color, body in ctas:
        add_rect(slide, left=cx, top=Inches(2.3), width=Inches(4.0), height=Inches(3.5),
                 fill_color=RGBColor(0x16, 0x1B, 0x22), line_color=color, line_width=Pt(2))
        add_rect(slide, left=cx, top=Inches(2.3), width=Inches(4.0), height=Inches(0.07),
                 fill_color=color)
        add_text(slide, title, left=cx + Inches(0.15), top=Inches(2.45),
                 width=Inches(3.7), height=Inches(0.6),
                 font_size=Pt(18), bold=True, color=color)
        add_text(slide, body, left=cx + Inches(0.15), top=Inches(3.1),
                 width=Inches(3.7), height=Inches(2.5),
                 font_size=Pt(15), color=WHITE)
        cx += Inches(4.3)

    # Quote at bottom
    add_text(slide, '"The programs are live. The SDK is published. Devnet is open."',
             left=Inches(0.5), top=Inches(5.95), width=Inches(12.3), height=Inches(0.5),
             font_size=Pt(17), color=GREEN, italic=True, align=PP_ALIGN.CENTER)

    set_notes(slide, """TIME: 4:15–5:00

"Three things I want you to do in the next 48 hours."

"One: install the SDK. npm install @holdfastprotocol/sdk@devnet. Connects to devnet. No mainnet risk. Register an agent identity. Query a reputation score. Tell us what's broken."

"Two: if you're building an agent framework or a DeFi protocol that uses agents, come find us at the conference. We're actively looking for integration partners before mainnet. That's an opportunity to shape the standard."

"Three: follow @CasemateLabs for the audit timeline and mainnet launch. The audit is in progress. When it's done, this goes to mainnet. You want to be building on it before that happens."

"The programs are live. The SDK is published. Devnet is open. The accountability layer for the agent economy — let's build it together."

END OF TALK — hold this slide until the clock hits 5:00.""")
    return slide


def slide_14_qa(prs):
    """Slide 14 — Q&A / Contact"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)

    # Green left bar
    add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5),
             fill_color=GREEN)

    # Logo / wordmark area
    add_text(slide, "Holdfast Protocol",
             left=Inches(0.5), top=Inches(1.2), width=Inches(12), height=Inches(1.1),
             font_size=Pt(58), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Trust infrastructure for the AI agent economy.",
             left=Inches(0.5), top=Inches(2.3), width=Inches(12), height=Inches(0.55),
             font_size=Pt(22), color=GREEN, align=PP_ALIGN.CENTER)

    add_accent_line(slide, color=GREEN, top=Inches(3.05))

    # Contact info
    contacts = [
        "Casemate Labs  ·  holdfastprotocol.com",
        "@CasemateLabs",
        "npm: @holdfastprotocol/sdk@devnet",
    ]
    y = Inches(3.25)
    for c in contacts:
        add_text(slide, c,
                 left=Inches(0.5), top=y, width=Inches(12), height=Inches(0.45),
                 font_size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.5)

    add_text(slide, "Questions?",
             left=Inches(0.5), top=Inches(4.75), width=Inches(12), height=Inches(0.65),
             font_size=Pt(32), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Disclaimer
    add_rect(slide, left=Inches(1.5), top=Inches(5.7), width=Inches(10.3), height=Inches(1.1),
             fill_color=RGBColor(0x2A, 0x16, 0x10), line_color=RED_WARN, line_width=Pt(1))
    add_text(slide, "⚠  Pre-audit software. Devnet only.\nNot for mainnet or production use.\nExternal security audit in progress.",
             left=Inches(1.7), top=Inches(5.78), width=Inches(9.9), height=Inches(0.95),
             font_size=Pt(13), color=RED_WARN, align=PP_ALIGN.CENTER)

    set_notes(slide, """Q&A ANTICIPATED QUESTIONS:

"How is this different from Vouch?"
→ "Vouch is an off-chain identity tool — signs Git commits, issues DIDs. Zero on-chain enforcement. No escrow, no dispute resolution, no financial consequences. We're the enforcement layer. They're actually a complementary integration partner."

"Why not use EAS on Ethereum?"
→ "EAS is general-purpose attestation on EVM. Not designed for AI agents, not on Solana, no escrow or financial enforcement. EVM gas at scale is a real constraint."

"When is mainnet?"
→ "After the external audit completes. We're not setting a hard date — audit timeline drives it. We'd rather ship secure than fast."

"Is there a token?"
→ "No token. Revenue model is protocol fees on real usage — registrations, escrow settlements."

"Can I use this with ElizaOS / Solana Agent Kit?"
→ "Integration proposals in progress. Come find us at the conference if you're building on those frameworks." """)
    return slide


# ===========================================================================
# MAIN
# ===========================================================================

def build():
    prs = new_prs()

    print("Building Holdfast Protocol pitch deck...")
    slide_01_title(prs);      print("  [ok] Slide 1  - Title")
    slide_02_problem(prs);    print("  [ok] Slide 2  - The Problem")
    slide_03_trust_gap(prs);  print("  [ok] Slide 3  - The Trust Gap")
    slide_04_intro(prs);      print("  [ok] Slide 4  - Introducing Holdfast Protocol")
    slide_05_identity(prs);   print("  [ok] Slide 5  - Primitive 1: Identity")
    slide_06_reputation(prs); print("  [ok] Slide 6  - Primitive 2: Reputation")
    slide_07_escrow(prs);     print("  [ok] Slide 7  - Primitive 3: Escrow")
    slide_08_why_solana(prs); print("  [ok] Slide 8  - Why Solana")
    slide_09_demo(prs);       print("  [ok] Slide 9  - Live Demo")
    slide_10_current_state(prs); print("  [ok] Slide 10 - Current State")
    slide_11_competitive(prs);   print("  [ok] Slide 11 - Competitive Differentiation")
    slide_12_ecosystem(prs);     print("  [ok] Slide 12 - Ecosystem Fit")
    slide_13_cta(prs);           print("  [ok] Slide 13 - Call to Action")
    slide_14_qa(prs);            print("  [ok] Slide 14 - Q&A / Contact")

    out = "docs/holdfast-pitch-deck-aiagentconf-2026.pptx"
    prs.save(out)
    print(f"\nDeck saved: {out}")
    print(f"   Slides: 14  |  Format: PPTX (widescreen 16:9)")
    print(f"   Upload to Google Slides: File > Import slides > Upload")
    print(f"   Or open directly in PowerPoint.")

if __name__ == "__main__":
    build()
